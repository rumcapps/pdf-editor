# app.py
from flask import Flask, request, jsonify, send_file, render_template
from flask_cors import CORS
from pdf_text_extractor import extract_text_from_pdf
from pptx import Presentation
from io import BytesIO

app = Flask(__name__)
CORS(app)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/convert', methods=['POST'])
def convert_pdf_to_pptx():
    try:
        pdf_file = request.files['pdf']
        pdf_path = 'uploaded_file.pdf'
        pdf_file.save(pdf_path)
        extracted_text = extract_text_from_pdf(pdf_path)

        # Placeholder logic to convert extracted text to PowerPoint
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title = slide.shapes.title
        title.text = extracted_text

        pptx_output = BytesIO()
        presentation.save(pptx_output)
        pptx_output.seek(0)

        # Remove the local PDF file after extraction
        # Uncomment the line below if you want to delete the PDF file
        # os.remove(pdf_path)

        return send_file(pptx_output, as_attachment=True, download_name='output.pptx')

    except Exception as e:
        error_message = str(e)
        app.logger.error(f"Error during conversion: {error_message}")
        response_data = {
            'success': False,
            'error': error_message,
        }
        return jsonify(response_data)

if __name__ == '__main__':
    app.run(debug=True)
